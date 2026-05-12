#!/usr/bin/env python3
"""
build_module_pptx.py
────────────────────
Build a module's PowerPoint deck from two inputs:

  1. lesson .md            slide content, talking points, tables
  2. <lesson>_image_prompts.md   sidecar with per-slide image prompts (optional)

Parsing
───────
Slides are extracted from the lesson by calling a small LLM (gpt-4o-mini) with
JSON output — not by regex. Each slide returns title, body_content, talking_points,
and table_md. The result is cached at <lesson>.parsed.json keyed by the lesson's
content hash; re-runs skip the LLM call when the lesson hasn't changed.

A `--regex-parser` flag falls back to the old regex parser if the LLM call fails
or you want a deterministic offline run.

Streaming dispatch
──────────────────
As each slide is parsed and classified, its image generation fires immediately
on a thread pool (MAX_WORKERS=50). The parser keeps moving while generations
run in the background. Once everything is dispatched, the script awaits all
in-flight generations, then assembles the .pptx.

Slide treatment is determined per slide:

  · prose prompt in sidecar         → OpenAI image generation
  · [existing: <path>] in sidecar   → use that image as-is, no generation
  · [text-only] in sidecar          → auto-template a prompt from the slide's
                                      body content and generate the image
  · lesson has a markdown table     → markdown → PDF → PNG → OpenAI stylize
  · (no spec, no table)             → ERROR (warn + white placeholder slide)

EVERY slide ends up as a full-bleed image. There are no native PPT text slides.

Output
──────
  <module-dir>/images/<slug>.png    slug derived from slide title
  <module-dir>/<lesson-stem>.pptx

Image filenames are slug-based — reordering slides does NOT rename images.
Title collisions get -2, -3 suffixes.

Usage
─────
  python3 build_module_pptx.py <lesson.md> [--slides N] [--no-images] [--regex-parser]
"""

import argparse
import base64
import hashlib
import json
import os
import re
import sys
import unicodedata
from concurrent.futures import ThreadPoolExecutor, as_completed
from dataclasses import dataclass, field
from pathlib import Path
from typing import Optional

from dotenv import load_dotenv
load_dotenv(Path(__file__).parent / ".env")

from openai import OpenAI
from lxml import etree
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn

from test_table_render import render_table_markdown_to_png
from test_table_stylize import stylize_table


# ── Constants ────────────────────────────────────────────────────────────────

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

C_BLACK   = RGBColor(0x00, 0x00, 0x00)
C_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
C_GRAY_LT = RGBColor(0xEB, 0xEB, 0xEB)
C_GRAY_MD = RGBColor(0x55, 0x55, 0x55)

IMAGE_SIZE          = "1536x1024"
IMAGE_MODEL         = "gpt-image-2"
TABLE_STYLIZE_SIZE  = "1536x1024"
PARSER_MODEL        = "gpt-4o-mini"     # small/cheap; only used for lesson parsing
MAX_WORKERS         = 50


# ── Data structures ──────────────────────────────────────────────────────────

@dataclass
class SlideSpec:
    kind: str                            # 'image' | 'existing' | 'text-only'
    prompt: Optional[str] = None
    existing_path: Optional[Path] = None


@dataclass
class SlideData:
    index: int
    title: str
    slug: str
    body_content: str = ""               # raw markdown body (preserves bold, lists, indent)
    body_lines: list = field(default_factory=list)
    table_rows: list = field(default_factory=list)
    table_md: Optional[str] = None
    talking_points: list = field(default_factory=list)      # legacy: stripped list of lines
    talking_points_raw: str = ""         # raw markdown talking-points (preserves indent)
    image_path: Optional[Path] = None
    spec: Optional[SlideSpec] = None
    treatment: str = "unspecified"


# ── Slug derivation ──────────────────────────────────────────────────────────

def slugify(text: str) -> str:
    text = re.sub(r'\s*\*\([^)]+\)\*', '', text)
    text = re.sub(r'\*\*(.+?)\*\*', r'\1', text)
    text = unicodedata.normalize('NFKD', text).encode('ascii', 'ignore').decode()
    text = re.sub(r'[^a-z0-9]+', '-', text.lower()).strip('-')
    return text or "slide"


def _clean_title(raw: str) -> str:
    return re.sub(r'\s*\*\([^)]+\)\*', '', raw).strip()


def _strip_md_bold(text: str) -> str:
    text = re.sub(r'\*\*(.+?)\*\*', r'\1', text)
    text = re.sub(r'\*(.+?)\*', r'\1', text)
    return text


# ── LLM-based parsing ────────────────────────────────────────────────────────

PARSER_SYSTEM = """You parse ISYS 398U course lesson markdown into structured slide data.

The lesson is a markdown document. Slides are marked by `### SLIDE: <title>` headings.
Some slides have a speaker-notes subsection marked by `#### 🗣 ...` (the heading may
say "Talking points", "Closing notes", "Facilitation", "Instructor notes", or similar —
treat all of these as speaker notes).

A lesson may also contain doc-level sections (single `##` headers, like
`## Learning objectives`) that are NOT marked with `### SLIDE:`. Treat any such
section that looks like presentation-worthy slide content (Learning Objectives,
Key Takeaways, etc.) as if it WERE a slide — promote it to a slide entry with
the section heading as its title.

For each slide, extract:

  title           — exact text after `### SLIDE: ` (or the `##` heading if promoted).
                    Preserve every character including em-dashes, quotes, parentheticals.
  body_content    — on-slide markdown content (everything between the slide heading
                    and either the speaker-notes subsection or the next slide). Keep
                    markdown formatting intact (**bold**, lists, tables, blockquotes).
  talking_points  — speaker-notes content (everything in the `#### 🗣 ...` subsection
                    until the next slide). Keep markdown intact. Empty string if none.
  table_md        — if body_content includes a markdown table (lines starting with `|`),
                    return JUST the table block as a string. Otherwise null.

Return JSON with this exact shape:

{
  "slides": [
    {"title": "...", "body_content": "...", "talking_points": "...", "table_md": "..." | null},
    ...
  ]
}

Preserve slide order. Do not invent slides. Do not modify content."""


def parse_lesson_llm(md_path: Path, cache_path: Optional[Path] = None) -> list[dict]:
    """Parse a lesson .md via LLM into a list of slide dicts. Cached by content hash."""
    text = md_path.read_text(encoding='utf-8')
    h = hashlib.sha256(text.encode()).hexdigest()[:16]

    if cache_path and cache_path.exists():
        try:
            cached = json.loads(cache_path.read_text())
            if cached.get("content_hash") == h:
                print(f"  [cached] LLM parse  ({len(cached['slides'])} slides)")
                return cached["slides"]
        except (json.JSONDecodeError, KeyError):
            pass

    print(f"  [llm]    parsing lesson via {PARSER_MODEL}…")
    client = OpenAI()
    response = client.chat.completions.create(
        model=PARSER_MODEL,
        messages=[
            {"role": "system", "content": PARSER_SYSTEM},
            {"role": "user", "content": f"Parse this lesson:\n\n{text}"},
        ],
        response_format={"type": "json_object"},
        temperature=0,
    )
    payload = json.loads(response.choices[0].message.content)
    slides = payload.get("slides", [])
    print(f"  [done]   parsed {len(slides)} slides")

    if cache_path:
        cache_path.write_text(json.dumps(
            {"content_hash": h, "slides": slides}, indent=2
        ))

    return slides


# ── Regex parser (fallback / --regex-parser flag) ────────────────────────────

def _parse_table_rows(section_text: str) -> list:
    rows = []
    for line in section_text.splitlines():
        if not re.match(r'^\s*\|', line):
            continue
        if re.match(r'^\s*\|[\-: |]+\|\s*$', line):
            continue
        cells = [_strip_md_bold(c.strip()) for c in line.split('|') if c.strip()]
        if cells:
            rows.append(cells)
    return rows


def _extract_table_md(section_text: str) -> Optional[str]:
    lines = []
    in_table = False
    for line in section_text.splitlines():
        if re.match(r'^\s*\|', line):
            lines.append(line)
            in_table = True
        elif in_table:
            break
    return "\n".join(lines) if lines else None


def parse_lesson_regex(md_path: Path) -> list[dict]:
    text = md_path.read_text(encoding='utf-8')
    slide_re = re.compile(r'^### SLIDE:\s*(.+)$', re.MULTILINE)
    matches = list(slide_re.finditer(text))
    out = []
    for i, m in enumerate(matches):
        title = m.group(1).strip()
        start = m.end()
        end = matches[i + 1].start() if i + 1 < len(matches) else len(text)
        section = text[start:end]
        tp_m = re.search(r'####\s+🗣[^\n]*\n(.*?)(?=\n---|\Z)', section, re.DOTALL)
        tp_text = tp_m.group(1).strip() if tp_m else ""
        body_text = section[:tp_m.start()].strip() if tp_m else section.strip()
        body_text = re.sub(r'\*\*\[Image prompt:.*?\]\*\*', '', body_text, flags=re.DOTALL)
        out.append({
            "title": title,
            "body_content": body_text,
            "talking_points": tp_text,
            "table_md": _extract_table_md(body_text),
        })
    return out


# ── SlideData construction from parser output ────────────────────────────────

def _parse_body_lines(text: str) -> list:
    """Returns list of (kind, text) — kind is 'bullet' or 'text'."""
    lines = []
    for line in text.splitlines():
        line = line.strip()
        if not line or line.startswith('#') or line == '---':
            continue
        if re.match(r'^\|', line):
            continue
        cleaned = _strip_md_bold(line)
        if cleaned.startswith('- '):
            lines.append(('bullet', cleaned[2:]))
        elif cleaned:
            lines.append(('text', cleaned))
    return lines


def build_slides(raw_slides: list[dict]) -> list[SlideData]:
    """Convert parser output dicts into SlideData with slugs + derived fields."""
    slides = []
    seen_counts = {}
    for i, raw in enumerate(raw_slides):
        title = raw["title"].strip()
        base = slugify(title)
        seen_counts[base] = seen_counts.get(base, 0) + 1
        slug = base if seen_counts[base] == 1 else f"{base}-{seen_counts[base]}"

        slide = SlideData(
            index=i + 1,
            title=title,
            slug=slug,
            body_content=raw.get("body_content", "") or "",
            table_md=raw.get("table_md") or None,
        )
        slide.body_lines = _parse_body_lines(slide.body_content)
        slide.table_rows = _parse_table_rows(slide.body_content)
        tp_text = raw.get("talking_points", "") or ""
        slide.talking_points_raw = tp_text
        slide.talking_points = [
            l.strip() for l in tp_text.splitlines() if l.strip()
        ]
        slides.append(slide)
    return slides


# ── Sidecar prompts file parsing ──────────────────────────────────────────────

def parse_prompts_file(path: Path) -> dict:
    if not path.exists():
        return {}
    text = path.read_text(encoding='utf-8')
    slide_re = re.compile(r'^##\s+SLIDE:\s*(.+)$', re.MULTILINE)
    matches = list(slide_re.finditer(text))
    specs = {}
    for i, m in enumerate(matches):
        title = m.group(1).strip()
        start = m.end()
        end = matches[i + 1].start() if i + 1 < len(matches) else len(text)
        body = text[start:end].strip()
        body = re.sub(r'\n\s*---\s*$', '', body).strip()
        if not body:
            continue
        if re.match(r'^\[text-only\]\s*$', body, re.IGNORECASE):
            specs[title] = SlideSpec(kind='text-only')
            continue
        ex_m = re.match(r'^\[existing:\s*([^\]]+)\]\s*$', body, re.IGNORECASE)
        if ex_m:
            specs[title] = SlideSpec(kind='existing', existing_path=Path(ex_m.group(1).strip()))
            continue
        specs[title] = SlideSpec(kind='image', prompt=body)
    return specs


def attach_specs(slides: list, prompts: dict) -> None:
    for slide in slides:
        spec = prompts.get(slide.title)
        if spec:
            slide.spec = spec
            slide.treatment = spec.kind
        elif slide.table_rows:
            slide.treatment = 'table'
        else:
            slide.treatment = 'unspecified'


# ── Auto-template for text-only slides (now image-generated) ─────────────────

def auto_prompt_for_text_only(slide: SlideData) -> str:
    """Build a colored-pencil-sketch prompt that renders the slide's body
    content as hand-lettered text. Used for slides marked [text-only] in the
    sidecar prompts file — they become full-bleed images like every other slide."""
    title = _clean_title(slide.title)
    # Use the raw body content so numbered lists and formatting survive.
    body = slide.body_content.strip()
    # Trim any leading/trailing markdown noise that shouldn't appear in art.
    body = _strip_md_bold(body)
    return (
        "Widescreen 16:9, sketched with colored pencils on a solid white background. "
        "A clean sheet of paper rendered in casual pencil style, with all text "
        "hand-lettered in the same colored-pencil aesthetic — no typeset fonts, "
        "no print typography. "
        f"At the top of the page, the title hand-lettered prominently: \"{title}\". "
        "Below the title, the following list rendered as hand-lettered numbered "
        "items, each kept on its own visual line, preserving the exact wording:\n\n"
        f"{body}\n\n"
        "Composition: clean and orderly, list-format, the page is the focal element. "
        "Muted earth-tone palette with one accent color. Mood: organized, "
        "instructive, the kind of page you'd photograph from a notebook."
    )


# ── Image generation — streaming dispatch ─────────────────────────────────────

def _gen_image(client: OpenAI, slide: SlideData, images_dir: Path, prompt: str) -> None:
    out = images_dir / f"{slide.slug}.png"
    if out.exists():
        print(f"  [cached] {slide.slug} → {out.name}")
        slide.image_path = out
        return
    print(f"  [gen]    {slide.slug} — requesting…")
    try:
        result = client.images.generate(
            model=IMAGE_MODEL,
            prompt=prompt,
            size=IMAGE_SIZE,
        )
        image_bytes = base64.b64decode(result.data[0].b64_json)
        out.write_bytes(image_bytes)
        slide.image_path = out
        print(f"  [done]   {slide.slug}")
    except Exception as exc:
        print(f"  [error]  {slide.slug}: {exc}", file=sys.stderr)


def _gen_table(client: OpenAI, slide: SlideData, images_dir: Path) -> None:
    out = images_dir / f"{slide.slug}.png"
    if out.exists():
        print(f"  [cached] {slide.slug} (table) → {out.name}")
        slide.image_path = out
        return

    title = _clean_title(slide.title)
    md_content = f"## {title}\n\n{slide.table_md}\n"
    structural = images_dir / f"{slide.slug}-structural.png"

    print(f"  [render] {slide.slug} (table) — pandoc→pdf→png…")
    try:
        render_table_markdown_to_png(md_content, structural)
    except Exception as exc:
        print(f"  [error]  {slide.slug} render: {exc}", file=sys.stderr)
        return

    print(f"  [gen]    {slide.slug} (table) — stylizing…")
    try:
        stylize_table(structural, out,
                      model=IMAGE_MODEL,
                      size=TABLE_STYLIZE_SIZE,
                      client=client)
        slide.image_path = out
        print(f"  [done]   {slide.slug} (table)")
    except Exception as exc:
        print(f"  [error]  {slide.slug} stylize: {exc}", file=sys.stderr)


def _resolve_existing(slide: SlideData, module_dir: Path) -> None:
    p = slide.spec.existing_path
    if not p.is_absolute():
        p = module_dir / p
    if p.exists():
        slide.image_path = p
        print(f"  [exist]  {slide.slug} → {p.name}")
    else:
        print(f"  [warn]   {slide.slug} existing path missing: {p}", file=sys.stderr)


def dispatch_all(slides: list, images_dir: Path, module_dir: Path,
                 skip_generation: bool = False) -> None:
    """Stream slides through dispatch. Every slide either resolves to an existing
    image, gets an image generated (prose / text-only-auto / table), or — in
    --no-images mode — picks up a cached image if one exists on disk."""
    images_dir.mkdir(exist_ok=True)
    client = None
    if not skip_generation:
        if not os.environ.get("OPENAI_API_KEY"):
            sys.exit("❌  OPENAI_API_KEY not set — add to Application/ImageGen/.env")
        client = OpenAI()

    with ThreadPoolExecutor(max_workers=MAX_WORKERS) as pool:
        futures = []
        for slide in slides:
            if slide.treatment == 'existing':
                _resolve_existing(slide, module_dir)
                continue

            # If a cached image already exists at the slug, use it regardless of mode
            cached = images_dir / f"{slide.slug}.png"

            if skip_generation:
                if cached.exists():
                    slide.image_path = cached
                continue

            if slide.treatment == 'image':
                futures.append(pool.submit(_gen_image, client, slide, images_dir, slide.spec.prompt))
            elif slide.treatment == 'text-only':
                # Auto-build a prompt from the slide's body content and dispatch
                prompt = auto_prompt_for_text_only(slide)
                futures.append(pool.submit(_gen_image, client, slide, images_dir, prompt))
            elif slide.treatment == 'table':
                futures.append(pool.submit(_gen_table, client, slide, images_dir))
            elif slide.treatment == 'unspecified':
                print(f"  [warn]   {slide.slug} — no spec, no table; will render as white placeholder",
                      file=sys.stderr)

        for f in as_completed(futures):
            f.result()


# ── PPTX helpers ─────────────────────────────────────────────────────────────

def _solid_rect(slide, left, top, width, height, color: RGBColor):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def _textbox(slide, left, top, width, height,
             text, font_size, color: RGBColor,
             bold=False, italic=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    return box


def _set_bg(slide_obj, color: RGBColor):
    fill = slide_obj.background.fill
    fill.solid()
    fill.fore_color.rgb = color


# ── Notes formatting ─────────────────────────────────────────────────────────

# Bullet characters per indent level
BULLET_CHARS = ['•', '◦', '▪', '▫']

def _parse_md_inline(text: str) -> list:
    """Parse **bold** and *italic* inline markdown. Returns (segment, bold) tuples.
    Italics are dropped to bold=False (they'd add formatting complexity in notes)."""
    result = []
    pos = 0
    for m in re.finditer(r'\*\*(.+?)\*\*', text):
        if m.start() > pos:
            result.append((text[pos:m.start()], False))
        result.append((m.group(1), True))
        pos = m.end()
    if pos < len(text):
        result.append((text[pos:], False))
    # Strip remaining single-* italic markers from each segment
    cleaned = []
    for seg, bold in result:
        cleaned.append((re.sub(r'\*([^*]+)\*', r'\1', seg), bold))
    return cleaned if cleaned else [(text, False)]


def _notes_para(txBody, segments, *, font_size_pt: int = 12,
                bold: bool = False, indent_level: int = 0,
                is_bullet: bool = False):
    """Append a paragraph to a notes txBody.

    font_size_pt — overall size for the paragraph's runs
    bold         — make ALL runs bold (overrides segment-level bold)
    indent_level — 0 = top level, 1 = sub, 2 = sub-sub, ...
    is_bullet    — render with a bullet character matching the indent level
    """
    p = etree.SubElement(txBody, qn('a:p'))
    pPr = etree.SubElement(p, qn('a:pPr'))
    if indent_level > 0:
        pPr.set('lvl', str(indent_level))
        # Manual indent (PPT notes don't always honor lvl visually)
        pPr.set('indent', str(-228600))               # hanging indent for bullet
        pPr.set('marL', str(457200 * (indent_level + 1)))  # left margin
    if is_bullet:
        char = BULLET_CHARS[min(indent_level, len(BULLET_CHARS) - 1)]
        buChar = etree.SubElement(pPr, qn('a:buChar'))
        buChar.set('char', char)
    for seg_text, seg_bold in segments:
        if not seg_text:
            continue
        r = etree.SubElement(p, qn('a:r'))
        rPr = etree.SubElement(r, qn('a:rPr'))
        rPr.set('lang', 'en-US')
        rPr.set('sz', str(font_size_pt * 100))         # PPT uses 1/100 pt
        if bold or seg_bold:
            rPr.set('b', '1')
        t = etree.SubElement(r, qn('a:t'))
        t.text = seg_text


def _render_md_to_notes(txBody, md_text: str) -> None:
    """Render a block of markdown into the notes txBody, preserving:
      - bullet hierarchy by leading whitespace (every 2 spaces = 1 indent level)
      - **bold** inline runs
      - numbered list items (rendered as text, the number stays in the line)
      - blank paragraphs between blocks
    Skips markdown headers, horizontal rules, and table syntax."""
    for raw_line in md_text.splitlines():
        if not raw_line.strip():
            etree.SubElement(txBody, qn('a:p'))
            continue
        stripped = raw_line.lstrip(' \t')
        # Skip markdown headers, horizontal rules, table rows
        if stripped.startswith('#'):
            continue
        if stripped.strip() == '---':
            continue
        if stripped.startswith('|'):
            continue

        indent_chars = len(raw_line) - len(stripped)
        indent_level = indent_chars // 2

        # Blockquote — italic prefix removed; render as quoted text
        if stripped.startswith('> '):
            inner = stripped[2:]
            _notes_para(txBody, _parse_md_inline(inner),
                        font_size_pt=12, indent_level=indent_level + 1)
            continue

        # Bullet
        if stripped.startswith('- ') or stripped.startswith('* '):
            text = stripped[2:]
            _notes_para(txBody, _parse_md_inline(text),
                        font_size_pt=12,
                        indent_level=indent_level, is_bullet=True)
            continue

        # Numbered list item — render as text (the number stays in the line)
        if re.match(r'^\d+\.\s', stripped):
            _notes_para(txBody, _parse_md_inline(stripped),
                        font_size_pt=12, indent_level=indent_level)
            continue

        # Plain paragraph (may have inline **bold**)
        _notes_para(txBody, _parse_md_inline(stripped),
                    font_size_pt=12, indent_level=indent_level)


def _set_notes(slide_obj, title: str, body_md: str, talking_points_md: str) -> None:
    """Build speaker notes from:
      - title — rendered as a large bold header
      - body_md — the slide's on-slide markdown (rendered under '— On slide —')
      - talking_points_md — speaker-notes markdown (rendered under '— Notes —')
    Each section preserves bullet hierarchy and bold inline formatting."""
    ntf = slide_obj.notes_slide.notes_text_frame
    txBody = ntf._txBody
    for p in list(txBody.findall(qn('a:p'))):
        txBody.remove(p)

    # Title — larger and bold
    if title:
        _notes_para(txBody, [(title, True)], font_size_pt=18, bold=True)
        etree.SubElement(txBody, qn('a:p'))   # blank line after title

    # — On slide — section
    if body_md.strip():
        _notes_para(txBody, [("— On slide —", True)], font_size_pt=13, bold=True)
        etree.SubElement(txBody, qn('a:p'))
        _render_md_to_notes(txBody, body_md)
        etree.SubElement(txBody, qn('a:p'))   # blank line after section

    # — Notes — section
    if talking_points_md.strip():
        _notes_para(txBody, [("— Notes —", True)], font_size_pt=13, bold=True)
        etree.SubElement(txBody, qn('a:p'))
        _render_md_to_notes(txBody, talking_points_md)


# ── Slide builder — every slide is a full-bleed image ────────────────────────

def build_image_slide(prs: Presentation, slide: SlideData):
    """Full-bleed image. Speaker notes include the slide title + the full body
    content + the talking points, with bullet hierarchy preserved. EVERY slide
    goes here."""
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    display_title = _clean_title(slide.title)
    if slide.image_path:
        _set_bg(sl, C_BLACK)
        sl.shapes.add_picture(str(slide.image_path), 0, 0, SLIDE_W, SLIDE_H)
    else:
        _set_bg(sl, C_WHITE)
    _set_notes(sl, display_title, slide.body_content, slide.talking_points_raw)
    return sl


# ── Deck assembly ────────────────────────────────────────────────────────────

def build_pptx(slides: list, out_path: Path) -> None:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    for slide in slides:
        build_image_slide(prs, slide)
    prs.save(str(out_path))
    print(f"\n✓  Saved → {out_path}")


# ── Entry point ──────────────────────────────────────────────────────────────

def main():
    parser = argparse.ArgumentParser(
        description="Stream-dispatch image generation for a module and assemble its .pptx."
    )
    parser.add_argument("markdown", help="Path to the module lesson .md file")
    parser.add_argument("--slides", "-n", type=int, default=None,
                        help="Max number of slides to process (omit for all)")
    parser.add_argument("--no-images", action="store_true",
                        help="Skip image generation; use cached images if present")
    parser.add_argument("--prompts", type=Path, default=None,
                        help="Path to image-prompts sidecar (default auto-detect)")
    parser.add_argument("--regex-parser", action="store_true",
                        help="Use the legacy regex parser instead of the LLM parser")
    args = parser.parse_args()

    md_path = Path(args.markdown).expanduser().resolve()
    if not md_path.exists():
        sys.exit(f"❌  File not found: {md_path}")

    if args.prompts:
        prompts_path = args.prompts.expanduser().resolve()
    else:
        stem = md_path.stem
        candidates = [
            md_path.with_name(f"{stem}_image_prompts.md"),
            md_path.with_name("image_prompts.md"),
            md_path.parent / f"{stem.split('_')[0]}_image_prompts.md",
        ]
        prompts_path = next((p for p in candidates if p.exists()), candidates[0])

    print(f"Lesson    {md_path.name}")
    print(f"Prompts   {prompts_path.name}{' (not found)' if not prompts_path.exists() else ''}")

    cache_path = md_path.with_suffix('.parsed.json')

    if args.regex_parser:
        print("Parser    regex (legacy)")
        raw_slides = parse_lesson_regex(md_path)
    else:
        print(f"Parser    LLM ({PARSER_MODEL})")
        try:
            raw_slides = parse_lesson_llm(md_path, cache_path)
        except Exception as exc:
            print(f"  [error]  LLM parse failed ({exc}); falling back to regex",
                  file=sys.stderr)
            raw_slides = parse_lesson_regex(md_path)

    if args.slides:
        raw_slides = raw_slides[:args.slides]

    slides = build_slides(raw_slides)
    prompts = parse_prompts_file(prompts_path)
    attach_specs(slides, prompts)

    counts = {}
    for s in slides:
        counts[s.treatment] = counts.get(s.treatment, 0) + 1
    print(f"Found {len(slides)} slides — " +
          ", ".join(f"{n} {k}" for k, n in counts.items()))
    print()

    module_dir = md_path.parent
    images_dir = module_dir / "images"

    if args.no_images:
        print("Images    → skipped (--no-images); cached images attached if present")
    else:
        print(f"Images    → {images_dir}")

    dispatch_all(slides, images_dir, module_dir, skip_generation=args.no_images)

    out_path = md_path.parent / (md_path.stem + ".pptx")
    print(f"\nBuilding  → {out_path}")
    build_pptx(slides, out_path)


if __name__ == "__main__":
    main()
