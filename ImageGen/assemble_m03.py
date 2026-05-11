#!/usr/bin/env python3
"""
assemble_m03.py
───────────────
Builds the complete M03 deck:
  [1]  Cover — M02 cover image (course visual identity)
  [2]  Module name slide — "Did the Agent Actually Do the Job?"
  [3–24] All 22 M03 content slides (images + tables)

Run from: Application/ImageGen/
  python3 assemble_m03.py
"""

import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

sys.path.insert(0, str(Path(__file__).parent.resolve()))
from build_module_pptx import (
    parse_markdown,
    build_image_slide,
    build_table_slide,
    _set_bg,
    _set_notes,
    _textbox,
    _solid_rect,
    C_BLACK,
    C_WHITE,
    C_GRAY_MD,
    SLIDE_W,
    SLIDE_H,
)

SCRIPT_DIR  = Path(__file__).parent.resolve()
REPO_ROOT   = SCRIPT_DIR.parent.parent

M02_IMAGES  = REPO_ROOT / "Module Content/M02 Inside an Agent's Head/images"
M03_DIR     = REPO_ROOT / "Module Content/M03 Grading the Agent: Evals and ROI"
M03_MD      = M03_DIR / "m03_grading_the_agent.md"
M03_IMAGES  = M03_DIR / "images"
M03_OUT     = M03_DIR / "m03_grading_the_agent.pptx"

MODULE_TITLE    = "Did the Agent Actually Do the Job?"
MODULE_SUBTITLE = "M03  ·  Evaluation & ROI  ·  ISYS 398U"


def build_cover_slide(prs):
    """Slide 1: Full-bleed M02 cover image — course visual identity."""
    cover_img = M02_IMAGES / "slide_01.png"
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    if cover_img.exists():
        _set_bg(sl, C_BLACK)
        sl.shapes.add_picture(str(cover_img), 0, 0, SLIDE_W, SLIDE_H)
        print(f"  ✓  Cover image  → {cover_img.name}")
    else:
        _set_bg(sl, C_WHITE)
        print(f"  ⚠  Cover image not found at {cover_img} — white placeholder used")
    _set_notes(sl, "ISYS 398U — Course Cover", [])
    return sl


def build_module_name_slide(prs):
    """Slide 2: Clean white slide — module title only, no image."""
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(sl, C_WHITE)

    # Thin rule
    _solid_rect(sl, Inches(0.55), Inches(2.75), Inches(12.2), Inches(0.012), C_BLACK)

    # Module title
    _textbox(sl,
             Inches(0.55), Inches(2.85),
             Inches(12.2), Inches(1.5),
             MODULE_TITLE,
             44, C_BLACK,
             bold=True,
             align=PP_ALIGN.LEFT)

    # Subtitle
    _textbox(sl,
             Inches(0.55), Inches(4.45),
             Inches(12.2), Inches(0.6),
             MODULE_SUBTITLE,
             18, C_GRAY_MD,
             bold=False,
             align=PP_ALIGN.LEFT)

    _set_notes(sl, MODULE_TITLE, [MODULE_SUBTITLE])
    return sl


def assign_images(slides):
    """Map each slide to its cached image file, if it exists."""
    for s in slides:
        if not s.image_prompt:
            s.image_path = None
            continue
        p = M03_IMAGES / f"slide_{s.index:02d}.png"
        s.image_path = p if p.exists() else None


def build_content_slide(prs, slide):
    if slide.table_rows:
        build_table_slide(prs, slide)
    else:
        build_image_slide(prs, slide)


def main():
    print(f"Module dir : {M03_DIR}")
    print(f"Images dir : {M03_IMAGES}\n")

    print(f"Parsing {M03_MD.name} …")
    slides = parse_markdown(M03_MD)
    print(f"  {len(slides)} slides\n")

    assign_images(slides)

    # Report image status
    print("Image status:")
    for s in slides:
        if s.image_prompt:
            status = f"✓ {s.image_path.name}" if s.image_path else "⚠  MISSING"
            print(f"  slide {s.index:02d}  {status}  — {s.title[:55]}")
        else:
            print(f"  slide {s.index:02d}  [table]              — {s.title[:55]}")

    missing = [s for s in slides if s.image_prompt and not s.image_path]
    if missing:
        print(f"\n⚠  {len(missing)} image(s) missing — those slides will render as white placeholders.")

    print("\nBuilding deck …")
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    print("  [1] Cover (M02 image) …")
    build_cover_slide(prs)

    print("  [2] Module name …")
    build_module_name_slide(prs)

    for i, slide in enumerate(slides):
        deck_num = i + 3
        kind = "table" if slide.table_rows else ("image" if slide.image_path else "placeholder")
        print(f"  [{deck_num:02d}] slide {slide.index:02d} — {kind:11s} — {slide.title[:50]}")
        build_content_slide(prs, slide)

    prs.save(str(M03_OUT))
    print(f"\n✓  Saved → {M03_OUT}")
    print(f"   Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
