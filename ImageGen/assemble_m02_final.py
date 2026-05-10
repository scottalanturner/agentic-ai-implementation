#!/usr/bin/env python3
"""
assemble_m02_final.py
─────────────────────
Builds the complete M02 deck by merging:
  - m02_inside_an_agents_head.md  (19 content slides)
  - m02_activities.md             (4 activity slides, inserted at the right spots)
  - A hand-built Learning Objectives slide (after the title)
  - Key takeaways already present as the last content slide

Also rebuilds m02_activities.pptx with real images.

Image note:
  Activities were generated to slots slide_01–slide_04.
  Main deck slides 1 (Title), 2 (Opener), and 4 (AI as labor)
  were cached from that earlier run — those three render as white
  placeholders in the merged deck. To fix them later:
    1. Delete images/slide_01.png, slide_02.png, slide_04.png from Finder
    2. Re-run: python3 build_module_pptx.py "../../Module Content/M02 Inside an Agent's Head/m02_inside_an_agents_head.md"
    3. Re-run this script

Run from: Application/ImageGen/
  python3 assemble_m02_final.py
"""

import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Paths ─────────────────────────────────────────────────────────────────────

SCRIPT_DIR = Path(__file__).parent.resolve()
REPO_ROOT  = SCRIPT_DIR.parent.parent
MODULE_DIR = REPO_ROOT / "Module Content/M02 Inside an Agent's Head"
IMAGES_DIR = MODULE_DIR / "images"
MAIN_MD    = MODULE_DIR / "m02_inside_an_agents_head.md"
ACT_MD     = MODULE_DIR / "m02_activities.md"
MERGED_OUT = MODULE_DIR / "m02_inside_an_agents_head.pptx"
ACT_OUT    = MODULE_DIR / "m02_activities.pptx"

# ── Import shared helpers from build_module_pptx ──────────────────────────────

sys.path.insert(0, str(SCRIPT_DIR))
from build_module_pptx import (
    parse_markdown,
    build_title_slide,
    build_image_slide,
    build_table_slide,
    _textbox,
    _solid_rect,
    _set_bg,
    _set_notes,
    _parse_md_inline,
    _notes_para,
    C_BLACK,
    C_WHITE,
    C_GRAY_LT,
    C_GRAY_MD,
    SLIDE_W,
    SLIDE_H,
)

# ── Learning objectives content ────────────────────────────────────────────────

LO_ITEMS = [
    ("Distinguish", "agents, chatbots, and traditional automation in realistic work scenarios"),
    ("Explain", "the ReAct-style loop — reason → act → observe — in plain English"),
    ("Explain", "what a context window is and why it shapes what an agent can know, remember, or do"),
    ("Name & illustrate", "four design patterns: reflection, tool use, planning, and multi-agent (preview)"),
    ("Distinguish", "system prompts from user prompts; explain why the Agent Card is the agent's SOP"),
    ("Identify", "human-in-the-loop (HITL) as a deliberate design choice, not an emergency brake"),
    ("Produce", "a first-pass Agent Card with a red team plan (P1) you can defend to a manager"),
]


def build_lo_slide(prs: Presentation):
    """Numbered learning-objectives slide — no title on the slide (title in notes).
    Items fill the full slide height for a clean, minimal look."""
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(sl, C_WHITE)

    ML   = Inches(0.55)
    TW   = Inches(12.2)

    # No title or rule on the slide — title goes to notes below
    ITEM_Y0 = Inches(0.45)   # start near top since title is removed
    ITEM_H  = Inches(0.88)   # more room per item

    for i, (verb, rest) in enumerate(LO_ITEMS):
        y = ITEM_Y0 + ITEM_H * i

        # Number + verb in one textbox (bold, medium-gray)
        num_box = sl.shapes.add_textbox(ML, y, Inches(1.55), ITEM_H)
        tf = num_box.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = f"{i + 1}.  {verb}"
        r.font.size      = Pt(14)
        r.font.bold      = True
        r.font.color.rgb = C_BLACK

        # Rest of the text (regular weight)
        txt_box = sl.shapes.add_textbox(ML + Inches(1.55), y,
                                        TW - Inches(1.55), ITEM_H)
        tf2 = txt_box.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.alignment = PP_ALIGN.LEFT
        r2 = p2.add_run()
        r2.text = rest
        r2.font.size      = Pt(14)
        r2.font.bold      = False
        r2.font.color.rgb = C_GRAY_MD

    # Title in notes (no text on slide)
    _set_notes(sl, "Learning Objectives", [])
    return sl


# ── Image assignment ──────────────────────────────────────────────────────────

# Main deck: which 1-based slide indices have wrong cached images?
# (activities were generated first, filling slots 01-04 before main deck ran)
MAIN_WRONG = {1, 2, 4}   # Title, Opener, AI as labor


def _assign_images(slides, wrong_indices=None):
    """Set slide.image_path from IMAGES_DIR/slide_NN.png.
    Slides in wrong_indices get None regardless of file existence."""
    wrong_indices = wrong_indices or set()
    for s in slides:
        if s.index in wrong_indices or not s.image_prompt:
            s.image_path = None
        else:
            p = IMAGES_DIR / f"slide_{s.index:02d}.png"
            s.image_path = p if p.exists() else None


# ── Generic slide dispatcher ──────────────────────────────────────────────────

def _build_slide(prs, slide):
    if slide.table_rows:
        build_table_slide(prs, slide)
    else:
        build_image_slide(prs, slide)


# ── Deck builders ─────────────────────────────────────────────────────────────

def build_merged(main_slides, act_slides, out_path: Path):
    """
    Final merged deck — 24 slides:

      [1]  Title
      [2]  Learning Objectives  ← new
      [3]  Opener
      [4]  Chatbot vs workflow vs agent (table)
      [5]  AI as labor
      [6]  The model alone is not an agent
      [7]  Context window
      [8]  ReAct in plain English
      [9]  Where loops go wrong
      [10] Human-in-the-loop (table)
      [11] Activity 1 — Trace the ReAct Loop
      [12] Why patterns matter
      [13] Pattern 1 — Reflection
      [14] Pattern 2 — Tool use
      [15] Pattern 3 — Planning
      [16] Pattern 4 — Multi-agent (teaser)
      [17] Activity 2 — Pattern Spotter
      [18] Two kinds of prompts (table)
      [19] A prompt is not "magic words"
      [20] Before / After
      [21] Activity 3 — SOP Makeover
      [22] Activity 4 — Failure Detective
      [23] P1 — What you're turning in
      [24] Key takeaways
    """
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    m = main_slides   # 0-based list; m[0] = main slide 1 (Title)
    a = act_slides    # 0-based list; a[0] = Activity 1

    # 1. Title
    build_title_slide(prs, m[0])
    # 2. Learning Objectives
    build_lo_slide(prs)
    # 3. Opener
    _build_slide(prs, m[1])
    # 4. Chatbot vs workflow vs agent
    _build_slide(prs, m[2])
    # 5. AI as labor
    _build_slide(prs, m[3])
    # 6. The model alone is not an agent
    _build_slide(prs, m[4])
    # 7. Context window
    _build_slide(prs, m[5])
    # 8. ReAct in plain English
    _build_slide(prs, m[6])
    # 9. Where loops go wrong
    _build_slide(prs, m[7])
    # 10. Human-in-the-loop
    _build_slide(prs, m[8])
    # 11. Activity 1 — Trace the ReAct Loop
    _build_slide(prs, a[0])
    # 12. Why patterns matter
    _build_slide(prs, m[9])
    # 13. Pattern 1 — Reflection
    _build_slide(prs, m[10])
    # 14. Pattern 2 — Tool use
    _build_slide(prs, m[11])
    # 15. Pattern 3 — Planning
    _build_slide(prs, m[12])
    # 16. Pattern 4 — Multi-agent
    _build_slide(prs, m[13])
    # 17. Activity 2 — Pattern Spotter
    _build_slide(prs, a[1])
    # 18. Two kinds of prompts
    _build_slide(prs, m[14])
    # 19. A prompt is not "magic words"
    _build_slide(prs, m[15])
    # 20. Before / After
    _build_slide(prs, m[16])
    # 21. Activity 3 — SOP Makeover
    _build_slide(prs, a[2])
    # 22. Activity 4 — Failure Detective
    _build_slide(prs, a[3])
    # 23. P1 Briefing
    _build_slide(prs, m[17])
    # 24. Key takeaways
    _build_slide(prs, m[18])

    prs.save(str(out_path))
    print(f"✓  Merged deck  → {out_path}  ({len(prs.slides)} slides)")


def build_activities_deck(act_slides, out_path: Path):
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    for s in act_slides:
        _build_slide(prs, s)
    prs.save(str(out_path))
    print(f"✓  Activities   → {out_path}  ({len(prs.slides)} slides)")


# ── Entry point ───────────────────────────────────────────────────────────────

def main():
    print(f"Module dir : {MODULE_DIR}")
    print(f"Images dir : {IMAGES_DIR}\n")

    # Parse
    print("Parsing m02_inside_an_agents_head.md …")
    main_slides = parse_markdown(MAIN_MD)
    print(f"  {len(main_slides)} slides\n")

    print("Parsing m02_activities.md …")
    act_slides = parse_markdown(ACT_MD)
    print(f"  {len(act_slides)} slides\n")

    # Assign images
    _assign_images(main_slides, wrong_indices=MAIN_WRONG)
    _assign_images(act_slides,  wrong_indices=set())

    # Report image status
    print("Image status (main deck):")
    for s in main_slides:
        status = "placeholder" if s.image_path is None else f"✓ {s.image_path.name}"
        if s.image_prompt:
            print(f"  slide {s.index:02d}  {status}")

    print("\nImage status (activities):")
    for s in act_slides:
        status = "placeholder" if s.image_path is None else f"✓ {s.image_path.name}"
        print(f"  activity {s.index}  {status}")

    # Build
    print("\nBuilding merged deck …")
    build_merged(main_slides, act_slides, MERGED_OUT)

    print("Building activities deck …")
    build_activities_deck(act_slides, ACT_OUT)

    print("\n━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━")
    print("⚠  Three main-deck images were cached from the activities run:")
    print("   slide_01.png  (Title)         — currently shows Activity 1 art")
    print("   slide_02.png  (Opener)        — currently shows Activity 2 art")
    print("   slide_04.png  (AI as labor)   — currently shows Activity 4 art")
    print()
    print("   To regenerate: delete those 3 files from Finder, then run:")
    print("   python3 build_module_pptx.py \\")
    print('     "../../Module Content/M02 Inside an Agent\'s Head/m02_inside_an_agents_head.md"')
    print("   Then re-run this script.")
    print("━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━")


if __name__ == "__main__":
    main()
